"""
PropChain — Prototype Progress PPT Generator
Team OpsAI | AI for Bharat Hackathon | Powered by AWS
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ─────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x1B, 0x3E)   # deep navy (background)
BLUE        = RGBColor(0x1A, 0x73, 0xE8)   # bright blue (headings)
TEAL        = RGBColor(0x00, 0xC2, 0xB2)   # teal (accent)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xE8, 0xF0, 0xFE)
GREEN       = RGBColor(0x34, 0xA8, 0x53)
AMBER       = RGBColor(0xFB, 0xBC, 0x04)
RED         = RGBColor(0xEA, 0x43, 0x35)
DARK_CARD   = RGBColor(0x16, 0x2A, 0x5A)   # card bg

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helpers ───────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank_layout)

def bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def txt(slide, text, l, t, w, h,
        size=20, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return tb

def badge(slide, label, l, t, color=TEAL):
    r = rect(slide, l, t, Inches(0.18), Inches(0.18), color)
    txt(slide, label, l + Inches(0.25), t - Inches(0.02),
        Inches(3.5), Inches(0.28), size=12, color=WHITE)

def divider(slide, t, color=TEAL):
    rect(slide, Inches(0.5), t, Inches(12.33), Inches(0.03), color)

def card(slide, l, t, w, h, title, lines, title_color=TEAL, icon=""):
    rect(slide, l, t, w, h, DARK_CARD)
    rect(slide, l, t, w, Inches(0.04), title_color)   # top accent bar
    txt(slide, f"{icon}  {title}" if icon else title,
        l + Inches(0.2), t + Inches(0.12), w - Inches(0.3), Inches(0.38),
        size=14, bold=True, color=title_color)
    body = "\n".join(lines)
    txt(slide, body,
        l + Inches(0.2), t + Inches(0.55), w - Inches(0.3), h - Inches(0.65),
        size=11, color=LIGHT_GRAY)

def pill(slide, text, l, t, color=BLUE):
    r = rect(slide, l, t, Inches(0.12), Inches(0.12), color)
    txt(slide, text, l + Inches(0.18), t - Inches(0.01), Inches(2), Inches(0.22),
        size=10, color=WHITE)
    return Inches(0.22)   # height consumed


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════
s = add_slide(); bg(s)

# left accent strip
rect(s, Inches(0), Inches(0), Inches(0.08), H, TEAL)

# big title
txt(s, "PropChain", Inches(0.5), Inches(1.6), Inches(7), Inches(1.6),
    size=72, bold=True, color=WHITE)
txt(s, "Blockchain & AI-Powered Property Registration",
    Inches(0.5), Inches(3.2), Inches(9), Inches(0.7),
    size=24, color=TEAL)

divider(s, Inches(4.1))

txt(s, "Prototype Progress — Phase 2 Submission",
    Inches(0.5), Inches(4.3), Inches(7), Inches(0.5),
    size=16, color=LIGHT_GRAY)
txt(s, "Team OpsAI  |  AI for Bharat Hackathon  |  Powered by AWS",
    Inches(0.5), Inches(4.9), Inches(8), Inches(0.4),
    size=13, color=RGBColor(0xAA, 0xBB, 0xDD))

# right side decoration
rect(s, Inches(10.2), Inches(1.5), Inches(2.6), Inches(4.5), DARK_CARD)
for i, (label, done) in enumerate([
    ("Mock Blockchain", True),
    ("AI Fraud Detection", True),
    ("Swagger UI + ReDoc", True),
    ("Postman Collection", True),
    ("React Frontend", False),
    ("Deployment", False),
]):
    col = GREEN if done else AMBER
    sym = "✓" if done else "○"
    txt(s, f"  {sym}  {label}",
        Inches(10.3), Inches(1.7) + Inches(0.62) * i,
        Inches(2.4), Inches(0.5),
        size=12, color=col)


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — What Was Built (Overview)
# ══════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, Inches(0), Inches(0), W, Inches(1.1), DARK_CARD)
txt(s, "What We Built", Inches(0.5), Inches(0.25), Inches(10), Inches(0.65),
    size=32, bold=True, color=WHITE)
txt(s, "Phase 2 Prototype — Backend complete",
    Inches(0.5), Inches(0.7), Inches(10), Inches(0.35),
    size=14, color=TEAL)

modules = [
    ("🔗", "Mock Blockchain", TEAL,
     ["SHA-256 hash-chained blocks", "MongoDB (append-only)", "6 transaction types",
      "Tamper detection", "Property Passport", "Block explorer"]),
    ("🤖", "AI Document Verification", BLUE,
     ["AWS Textract (extract)", "AWS Bedrock / Claude (analyze)", "7 rule-based checks",
      "fraud_score 0–1", "Auto on-chain logging", "Mock mode (no AWS needed)"]),
    ("📡", "REST API — 14 endpoints", RGBColor(0x9C, 0x27, 0xB0),
     ["POST /blockchain/mint", "POST /blockchain/transfer", "POST /ai/verify-document",
      "GET  /blockchain/passport", "GET  /blockchain/verify", "GET  /ai/mode"]),
    ("📖", "Docs & Tooling", RGBColor(0xE6, 0x7E, 0x22),
     ["Swagger UI  → /swagger", "ReDoc        → /redoc", "Postman collection (15 requests)",
      "run.sh command file", "Commands.md reference", ".env.example with AWS keys"]),
]

for i, (icon, title, col, lines) in enumerate(modules):
    c = i % 2
    r = i // 2
    card(s,
         Inches(0.4)  + Inches(6.3) * c,
         Inches(1.3)  + Inches(2.95) * r,
         Inches(6.1), Inches(2.8),
         title, lines, col, icon)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — Mock Blockchain Deep Dive
# ══════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, Inches(0), Inches(0), W, Inches(1.1), DARK_CARD)
txt(s, "Mock Blockchain", Inches(0.5), Inches(0.2), Inches(10), Inches(0.65),
    size=32, bold=True, color=WHITE)
txt(s, "SHA-256 hash-chained ledger backed by MongoDB",
    Inches(0.5), Inches(0.72), Inches(10), Inches(0.35), size=14, color=TEAL)

# Block structure diagram
rect(s, Inches(0.4), Inches(1.3), Inches(4.2), Inches(5.8), DARK_CARD)
txt(s, "Block Structure", Inches(0.6), Inches(1.4), Inches(3.8), Inches(0.4),
    size=15, bold=True, color=TEAL)
divider(s, Inches(1.85), TEAL)

fields = [
    ("block_index",      "0, 1, 2 …",           LIGHT_GRAY),
    ("property_id",      "PROP-KA-2024-001",     LIGHT_GRAY),
    ("transaction_type", "GENESIS | TRANSFER …", AMBER),
    ("data",             "{ owner, price, … }",  LIGHT_GRAY),
    ("previous_hash",    "000…0 (genesis)",      TEAL),
    ("timestamp",        "2026-02-28T18:00Z",    LIGHT_GRAY),
    ("hash",             "SHA-256( all above )", GREEN),
    ("minted_by",        "PropChain-Mock",       LIGHT_GRAY),
]
for i, (key, val, col) in enumerate(fields):
    y = Inches(1.95) + Inches(0.56) * i
    txt(s, key,  Inches(0.6),  y, Inches(1.6), Inches(0.45), size=11, bold=True, color=col)
    txt(s, val,  Inches(2.3),  y, Inches(2.1), Inches(0.45), size=10, color=LIGHT_GRAY)

# Transaction types
rect(s, Inches(5.0), Inches(1.3), Inches(7.9), Inches(2.55), DARK_CARD)
txt(s, "6 Transaction Types", Inches(5.2), Inches(1.4), Inches(7), Inches(0.4),
    size=15, bold=True, color=TEAL)
tx_types = [
    ("GENESIS",               "First property registration",          GREEN),
    ("OWNERSHIP_TRANSFER",    "Buyer ↔ Seller transfer",              BLUE),
    ("DOCUMENT_VERIFICATION", "AI fraud check result logged",         AMBER),
    ("STATUS_UPDATE",         "Legal status (disputed, encumbered)",  RED),
    ("FRACTIONAL_MINT",       "Property tokenized into shares",       TEAL),
    ("FRACTIONAL_TRANSFER",   "Token sold to investor",               RGBColor(0x9C,0x27,0xB0)),
]
for i, (name, desc, col) in enumerate(tx_types):
    y = Inches(1.85) + Inches(0.34) * i
    rect(s, Inches(5.1), y + Inches(0.06), Inches(0.1), Inches(0.18), col)
    txt(s, name, Inches(5.3), y, Inches(2.4), Inches(0.32), size=10, bold=True, color=col)
    txt(s, desc, Inches(7.75), y, Inches(4.9), Inches(0.32), size=10, color=LIGHT_GRAY)

# Security features
rect(s, Inches(5.0), Inches(4.05), Inches(7.9), Inches(3.05), DARK_CARD)
txt(s, "Security & Integrity", Inches(5.2), Inches(4.15), Inches(7), Inches(0.4),
    size=15, bold=True, color=TEAL)
sec = [
    (GREEN,  "Append-only",     "MongoDB unique index on (property_id, block_index)"),
    (GREEN,  "Tamper detection","Re-compute hash on every GET /verify → catch DB edits"),
    (GREEN,  "Chain linkage",   "Each block stores previous_hash → broken chain detected"),
    (GREEN,  "Genesis anchor",  "Block 0 previous_hash = 000…0 (64 zeros) — immutable"),
    (AMBER,  "Duplicate guard", "400 error if property already registered"),
]
for i, (col, title, desc) in enumerate(sec):
    y = Inches(4.6) + Inches(0.45) * i
    rect(s, Inches(5.1), y + Inches(0.08), Inches(0.1), Inches(0.18), col)
    txt(s, f"{title}:", Inches(5.3), y, Inches(2.2), Inches(0.38), size=11, bold=True, color=col)
    txt(s, desc, Inches(7.6), y, Inches(5.1), Inches(0.38), size=10, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — AI Document Verification
# ══════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, Inches(0), Inches(0), W, Inches(1.1), DARK_CARD)
txt(s, "AI Document Verification", Inches(0.5), Inches(0.2), Inches(10), Inches(0.65),
    size=32, bold=True, color=WHITE)
txt(s, "AWS Textract + Bedrock (Claude Haiku)  |  Mock fallback — no AWS needed",
    Inches(0.5), Inches(0.72), Inches(11), Inches(0.35), size=14, color=TEAL)

# Pipeline steps
steps = [
    (TEAL,  "1",  "Upload",        "PDF / JPEG / PNG\nMax 10 MB"),
    (BLUE,  "2",  "Hash",          "SHA-256 of file bytes\nDocument identity"),
    (AMBER, "3",  "Textract",      "Extract key-value\nfields from document"),
    (TEAL,  "4",  "Rule Checks",   "7 rule-based fraud\nchecks (no AWS)"),
    (BLUE,  "5",  "Bedrock/Claude","Holistic AI fraud\nanalysis + score"),
    (GREEN, "6",  "Verdict",       "AUTHENTIC / SUSPICIOUS\n/ FLAGGED"),
    (RGBColor(0x9C,0x27,0xB0), "7", "On-Chain", "Logged as\nDOC_VERIFICATION block"),
]

for i, (col, num, title, desc) in enumerate(steps):
    x = Inches(0.35) + Inches(1.85) * i
    rect(s, x, Inches(1.3), Inches(1.65), Inches(2.2), DARK_CARD)
    rect(s, x, Inches(1.3), Inches(1.65), Inches(0.05), col)
    # number circle
    c = rect(s, x + Inches(0.6), Inches(1.45), Inches(0.45), Inches(0.45), col)
    txt(s, num, x + Inches(0.6), Inches(1.44), Inches(0.45), Inches(0.45),
        size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, x + Inches(0.1), Inches(2.05), Inches(1.45), Inches(0.4),
        size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, desc, x + Inches(0.1), Inches(2.5), Inches(1.45), Inches(0.8),
        size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    # arrow
    if i < len(steps) - 1:
        txt(s, "→", x + Inches(1.65), Inches(1.85), Inches(0.25), Inches(0.4),
            size=18, color=TEAL, align=PP_ALIGN.CENTER)

# Rule checks grid
rect(s, Inches(0.35), Inches(3.75), Inches(12.6), Inches(3.4), DARK_CARD)
txt(s, "7 Rule-Based Fraud Checks  (run on every upload, no AWS required)",
    Inches(0.55), Inches(3.85), Inches(12), Inches(0.4),
    size=14, bold=True, color=TEAL)

rules = [
    ("Future Date Detection",         "Registration/execution dates must not be in the future"),
    ("Registration Number Format",    "Indian reg numbers match state-specific patterns"),
    ("Name Consistency",              "Owner name must be consistent across fields"),
    ("Aadhaar Format Validation",     "Aadhaar must be exactly 12 digits"),
    ("Amount Format Validation",      "Consideration amounts must be numeric"),
    ("Mandatory Fields Check",        "All required fields present for document type"),
    ("Extraction Confidence Check",   "Fields with Textract confidence < 70% are flagged"),
]

for i, (rule, desc) in enumerate(rules):
    col = i % 3
    row = i // 3
    x = Inches(0.5)  + Inches(4.2) * col
    y = Inches(4.35) + Inches(0.92) * row
    rect(s, x, y, Inches(4.0), Inches(0.82), RGBColor(0x0F, 0x22, 0x4A))
    rect(s, x, y, Inches(0.06), Inches(0.82), GREEN)
    txt(s, rule, x + Inches(0.18), y + Inches(0.05), Inches(3.7), Inches(0.32),
        size=11, bold=True, color=GREEN)
    txt(s, desc, x + Inches(0.18), y + Inches(0.38), Inches(3.7), Inches(0.38),
        size=9.5, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — Fraud Scoring & Verdict
# ══════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, Inches(0), Inches(0), W, Inches(1.1), DARK_CARD)
txt(s, "Fraud Scoring & Verdict Engine", Inches(0.5), Inches(0.2), Inches(12), Inches(0.65),
    size=32, bold=True, color=WHITE)
txt(s, "Composite score from AI analysis + rule penalties → deterministic verdict",
    Inches(0.5), Inches(0.72), Inches(12), Inches(0.35), size=14, color=TEAL)

# Score formula
rect(s, Inches(0.4), Inches(1.2), Inches(6.2), Inches(2.0), DARK_CARD)
txt(s, "Score Formula", Inches(0.6), Inches(1.3), Inches(5.5), Inches(0.38),
    size=15, bold=True, color=TEAL)
txt(s, "final_score = min(ai_score + (failed_rules × 0.08), 1.0)",
    Inches(0.6), Inches(1.75), Inches(5.8), Inches(0.38),
    size=13, bold=True, color=AMBER)
txt(s, "• AI score (0–1) from Bedrock/Claude analysis\n• Each failed rule adds 8% to the score\n• Capped at 1.0",
    Inches(0.6), Inches(2.2), Inches(5.8), Inches(0.85), size=11, color=LIGHT_GRAY)

# Verdict thresholds
rect(s, Inches(6.9), Inches(1.2), Inches(6.0), Inches(2.0), DARK_CARD)
txt(s, "Verdict Thresholds", Inches(7.1), Inches(1.3), Inches(5.5), Inches(0.38),
    size=15, bold=True, color=TEAL)
thresholds = [
    (GREEN, "AUTHENTIC",  "fraud_score 0.00 – 0.34",  "Document is clean"),
    (AMBER, "SUSPICIOUS", "fraud_score 0.35 – 0.65",  "Manual review needed"),
    (RED,   "FLAGGED",    "fraud_score 0.66 – 1.00",  "High fraud risk — block"),
]
for i, (col, verdict, score, action) in enumerate(thresholds):
    y = Inches(1.75) + Inches(0.42) * i
    rect(s, Inches(7.0), y + Inches(0.05), Inches(0.14), Inches(0.24), col)
    txt(s, verdict, Inches(7.25), y, Inches(1.5), Inches(0.35), size=12, bold=True, color=col)
    txt(s, score,   Inches(8.8),  y, Inches(2.2), Inches(0.35), size=11, color=LIGHT_GRAY)
    txt(s, action,  Inches(11.1), y, Inches(1.7), Inches(0.35), size=10, color=LIGHT_GRAY)

# Example verdicts
examples = [
    (GREEN,  "AUTHENTIC",  "0.04", "Title Deed — all fields valid, dates correct, name consistent"),
    (GREEN,  "AUTHENTIC",  "0.12", "Aadhaar Card — 12-digit number, valid DOB, high confidence"),
    (AMBER,  "SUSPICIOUS", "0.52", "Sale Agreement — missing stamp duty ref, low-confidence scan"),
    (RED,    "FLAGGED",    "0.87", "Title Deed — future registration date, name mismatch, bad reg no."),
    (RED,    "FLAGGED",    "1.00", "Aadhaar — future DOB (2030), only 8-digit number, address unknown"),
]

rect(s, Inches(0.4), Inches(3.35), Inches(12.5), Inches(3.8), DARK_CARD)
txt(s, "Live Examples (Mock Mode)", Inches(0.6), Inches(3.45), Inches(11), Inches(0.38),
    size=15, bold=True, color=TEAL)
divider(s, Inches(3.9), TEAL)

headers = ["Verdict", "Score", "Document & Finding"]
for hi, h in enumerate(headers):
    x = [Inches(0.6), Inches(2.4), Inches(3.7)][hi]
    txt(s, h, x, Inches(3.95), Inches(2), Inches(0.3), size=11, bold=True, color=TEAL)

for i, (col, verdict, score, desc) in enumerate(examples):
    y = Inches(4.35) + Inches(0.52) * i
    rect(s, Inches(0.5), y, Inches(12.3), Inches(0.46),
         RGBColor(0x0F, 0x22, 0x4A) if i % 2 == 0 else DARK_CARD)
    txt(s, verdict, Inches(0.6), y + Inches(0.08), Inches(1.7), Inches(0.32),
        size=11, bold=True, color=col)
    txt(s, score,   Inches(2.4), y + Inches(0.08), Inches(1.0), Inches(0.32),
        size=11, color=col)
    txt(s, desc,    Inches(3.7), y + Inches(0.08), Inches(9.0), Inches(0.32),
        size=10, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — API Endpoints
# ══════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, Inches(0), Inches(0), W, Inches(1.1), DARK_CARD)
txt(s, "API Endpoints — 14 Total", Inches(0.5), Inches(0.2), Inches(10), Inches(0.65),
    size=32, bold=True, color=WHITE)
txt(s, "FastAPI  |  Swagger UI → /swagger  |  ReDoc → /redoc  |  OpenAPI JSON → /openapi.json",
    Inches(0.5), Inches(0.72), Inches(12), Inches(0.35), size=13, color=TEAL)

# Left column — Blockchain
rect(s, Inches(0.3), Inches(1.2), Inches(6.3), Inches(5.9), DARK_CARD)
txt(s, "🔗  Blockchain", Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.4),
    size=16, bold=True, color=TEAL)
divider(s, Inches(1.78), TEAL)

blockchain_eps = [
    ("POST", "/blockchain/mint",                    "Register property (genesis block)",  GREEN),
    ("POST", "/blockchain/transfer",                "Record ownership transfer",           BLUE),
    ("POST", "/blockchain/verify-document",         "Log AI fraud result on-chain",        AMBER),
    ("POST", "/blockchain/status",                  "Update legal status",                 AMBER),
    ("POST", "/blockchain/fractional/mint",         "Tokenize property into shares",       TEAL),
    ("POST", "/blockchain/fractional/transfer",     "Transfer fractional tokens",          TEAL),
    ("GET",  "/blockchain/passport/{property_id}",  "Full Property Passport",              GREEN),
    ("GET",  "/blockchain/verify/{property_id}",    "Chain integrity check",               GREEN),
    ("GET",  "/blockchain/block/{block_hash}",      "Block explorer",                      BLUE),
    ("GET",  "/blockchain/properties",              "List all properties",                 BLUE),
]
for i, (method, path, desc, col) in enumerate(blockchain_eps):
    y = Inches(1.9) + Inches(0.48) * i
    mc = GREEN if method == "GET" else BLUE
    rect(s, Inches(0.4), y + Inches(0.05), Inches(0.55), Inches(0.28), mc)
    txt(s, method, Inches(0.4), y + Inches(0.03), Inches(0.55), Inches(0.3),
        size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, path, Inches(1.05), y, Inches(3.1), Inches(0.35), size=9.5, bold=True, color=LIGHT_GRAY)
    txt(s, desc, Inches(1.05), y + Inches(0.24), Inches(5.3), Inches(0.22), size=8.5, color=RGBColor(0x88,0x99,0xBB))

# Right column — AI + Health
rect(s, Inches(6.9), Inches(1.2), Inches(6.1), Inches(2.6), DARK_CARD)
txt(s, "🤖  AI Document Verification", Inches(7.1), Inches(1.3), Inches(5.6), Inches(0.4),
    size=16, bold=True, color=BLUE)
divider(s, Inches(1.78), BLUE)

ai_eps = [
    ("POST", "/ai/verify-document", "Upload PDF/image → full AI pipeline + on-chain log"),
    ("GET",  "/ai/mode",            "Check AWS or mock mode status"),
]
for i, (method, path, desc) in enumerate(ai_eps):
    y = Inches(1.9) + Inches(0.7) * i
    mc = GREEN if method == "GET" else BLUE
    rect(s, Inches(7.0), y + Inches(0.05), Inches(0.55), Inches(0.28), mc)
    txt(s, method, Inches(7.0), y + Inches(0.03), Inches(0.55), Inches(0.3),
        size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, path, Inches(7.65), y, Inches(5.1), Inches(0.35), size=10, bold=True, color=LIGHT_GRAY)
    txt(s, desc, Inches(7.65), y + Inches(0.35), Inches(5.1), Inches(0.28), size=9, color=RGBColor(0x88,0x99,0xBB))

rect(s, Inches(6.9), Inches(4.0), Inches(6.1), Inches(3.1), DARK_CARD)
txt(s, "📖  Docs & Tooling", Inches(7.1), Inches(4.1), Inches(5.5), Inches(0.4),
    size=16, bold=True, color=AMBER)
divider(s, Inches(4.58), AMBER)

tools = [
    ("Swagger UI",        "/swagger",         "Interactive API explorer — try it out"),
    ("ReDoc",             "/redoc",           "Clean read-only documentation"),
    ("OpenAPI JSON",      "/openapi.json",    "Machine-readable spec"),
    ("Postman Collection","15 requests",      "Auto-saved variables + test scripts"),
    ("run.sh",            "9 commands",       "setup / start / stop / status / open"),
    ("Commands.md",       "Full reference",   "All endpoints, env vars, structure"),
]
for i, (tool, val, desc) in enumerate(tools):
    y = Inches(4.65) + Inches(0.4) * i
    txt(s, tool, Inches(7.1), y, Inches(2.0), Inches(0.35), size=10, bold=True, color=AMBER)
    txt(s, val,  Inches(9.2), y, Inches(1.4), Inches(0.35), size=10, color=TEAL)
    txt(s, desc, Inches(10.7),y, Inches(2.1), Inches(0.35), size=9,  color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — AWS Integration
# ══════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, Inches(0), Inches(0), W, Inches(1.1), DARK_CARD)
txt(s, "AWS Integration", Inches(0.5), Inches(0.2), Inches(10), Inches(0.65),
    size=32, bold=True, color=WHITE)
txt(s, "Code is ready — plug in credentials to switch from mock → real AI",
    Inches(0.5), Inches(0.72), Inches(12), Inches(0.35), size=14, color=TEAL)

services = [
    (TEAL,  "Amazon Textract",
     "Extracts key-value pairs from property documents (PDF, JPEG, PNG, TIFF).\nUsed in: _textract_extract() + _textract_raw_text()",
     ["analyze_document (FORMS)", "detect_document_text", "Direct bytes (< 5MB)", "Returns field confidence scores"]),
    (BLUE,  "Amazon Bedrock — Claude Haiku",
     "Holistic fraud analysis. Receives extracted fields + raw text → returns fraud score, indicators, explanation.\nUsed in: _bedrock_analyze()",
     ["Model: anthropic.claude-haiku-4-5-20251001", "Structured JSON response", "Prompt engineered for Indian property docs", "Max 1024 tokens"]),
    (AMBER, "Auto-Switch Logic",
     "No code change needed. Server detects AWS credentials at startup and switches mode automatically.",
     ["USE_AWS = bool(AWS_ACCESS_KEY_ID)", "Mock mode if no credentials", "GET /ai/mode shows current mode", "Same API contract in both modes"]),
]

for i, (col, title, desc, bullets) in enumerate(services):
    y = Inches(1.3) + Inches(2.0) * i
    rect(s, Inches(0.4), y, Inches(12.5), Inches(1.85), DARK_CARD)
    rect(s, Inches(0.4), y, Inches(0.08), Inches(1.85), col)
    txt(s, title, Inches(0.7), y + Inches(0.1), Inches(4), Inches(0.42),
        size=16, bold=True, color=col)
    txt(s, desc, Inches(0.7), y + Inches(0.55), Inches(5.5), Inches(0.9),
        size=10.5, color=LIGHT_GRAY)
    for j, b in enumerate(bullets):
        bx = Inches(6.5) + Inches(3.1) * (j % 2)
        by = y + Inches(0.2) + Inches(0.42) * (j // 2)
        rect(s, bx, by + Inches(0.1), Inches(0.1), Inches(0.14), col)
        txt(s, b, bx + Inches(0.2), by, Inches(2.8), Inches(0.38), size=10, color=LIGHT_GRAY)

# To enable box
rect(s, Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.28), RGBColor(0x0F, 0x22, 0x4A))
txt(s, "To enable real AWS:  Add  AWS_ACCESS_KEY_ID  +  AWS_SECRET_ACCESS_KEY  +  AWS_REGION  to backend/.env  →  restart server",
    Inches(0.6), Inches(7.11), Inches(12.0), Inches(0.26),
    size=11, bold=True, color=AMBER)


# ══════════════════════════════════════════════════════════════
# SLIDE 8 — Tech Stack
# ══════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, Inches(0), Inches(0), W, Inches(1.1), DARK_CARD)
txt(s, "Tech Stack", Inches(0.5), Inches(0.2), Inches(10), Inches(0.65),
    size=32, bold=True, color=WHITE)
txt(s, "Current prototype  |  React Frontend & Deployment coming next",
    Inches(0.5), Inches(0.72), Inches(12), Inches(0.35), size=14, color=TEAL)

stack = [
    (TEAL,                       "Backend",    "FastAPI (Python 3.12)",        "Async, OpenAPI auto-docs, Pydantic validation"),
    (BLUE,                       "Database",   "MongoDB + Motor",              "Async driver, append-only blockchain ledger"),
    (GREEN,                      "Blockchain", "SHA-256 Mock Chain",           "Hash-chained blocks, tamper detection, 6 tx types"),
    (AMBER,                      "AI",         "AWS Textract + Bedrock",       "Document extraction + Claude Haiku fraud analysis"),
    (RED,                        "Fraud Rules","7 Rule-Based Checks",          "Python — no AWS, always runs, date/format/name checks"),
    (RGBColor(0x9C, 0x27, 0xB0), "API Docs",   "Swagger UI + ReDoc",           "FastAPI built-in, custom dark theme at /swagger"),
    (TEAL,                       "Tooling",    "Postman + run.sh",             "15-request collection, shell runner, Commands.md"),
    (BLUE,                       "Frontend",   "React.js  (planned)",          "Property dashboard, passport viewer, fraud results"),
    (GREEN,                      "Deployment", "AWS EC2 / Railway  (planned)", "Live URL for hackathon submission"),
]

for i, (col, layer, tech, desc) in enumerate(stack):
    row = i // 3
    col_n = i % 3
    x = Inches(0.35) + Inches(4.32) * col_n
    y = Inches(1.3)  + Inches(1.98) * row
    rect(s, x, y, Inches(4.1), Inches(1.82), DARK_CARD)
    rect(s, x, y, Inches(4.1), Inches(0.06), col)
    txt(s, layer, x + Inches(0.18), y + Inches(0.12), Inches(3.7), Inches(0.35),
        size=11, color=col)
    txt(s, tech,  x + Inches(0.18), y + Inches(0.48), Inches(3.7), Inches(0.42),
        size=14, bold=True, color=WHITE)
    txt(s, desc,  x + Inches(0.18), y + Inches(0.95), Inches(3.7), Inches(0.72),
        size=10, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 9 — What's Next
# ══════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, Inches(0), Inches(0), W, Inches(1.1), DARK_CARD)
txt(s, "What's Next", Inches(0.5), Inches(0.2), Inches(10), Inches(0.65),
    size=32, bold=True, color=WHITE)
txt(s, "Remaining work to complete the hackathon submission",
    Inches(0.5), Inches(0.72), Inches(12), Inches(0.35), size=14, color=TEAL)

next_items = [
    (BLUE,  "React Frontend",     "HIGH",
     ["Property registration form", "Document upload with live AI result",
      "Property Passport timeline view", "Fractional ownership marketplace"]),
    (AMBER, "AWS Credentials",    "HIGH",
     ["Add keys to backend/.env", "Enable Bedrock Claude Haiku in us-east-1",
      "Test with real property PDF", "GET /ai/mode → 'aws'"]),
    (GREEN, "Deployment",         "HIGH",
     ["AWS EC2 (t2.micro free tier)  or  Railway.app", "Live URL for evaluators",
      "MongoDB Atlas for production DB", "Update Postman base_url"]),
    (RED,   "GitHub Repository",  "MED",
     ["Clean commit history", "README with setup + demo GIF",
      ".gitignore (exclude .env)", "Architecture diagram"]),
    (TEAL,  "Demo Video",         "MED",
     ["Record 3-min screen walkthrough", "Show: Register → Verify (flagged) → Passport",
      "Show: Fractional investment flow", "Upload to YouTube / Loom"]),
    (RGBColor(0x9C,0x27,0xB0), "Project Summary", "LOW",
     ["Problem statement (from PPT)", "Solution impact numbers",
      "AWS services used", "Future roadmap"]),
]

for i, (col, title, priority, bullets) in enumerate(next_items):
    row = i // 3
    cn  = i % 3
    x = Inches(0.35) + Inches(4.32) * cn
    y = Inches(1.3)  + Inches(2.88) * row
    rect(s, x, y, Inches(4.1), Inches(2.72), DARK_CARD)
    rect(s, x, y, Inches(4.1), Inches(0.06), col)
    pc = RED if priority == "HIGH" else AMBER if priority == "MED" else GREEN
    rect(s, x + Inches(3.0), y + Inches(0.12), Inches(0.88), Inches(0.3), pc)
    txt(s, priority, x + Inches(3.0), y + Inches(0.12), Inches(0.88), Inches(0.3),
        size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, x + Inches(0.18), y + Inches(0.12), Inches(2.7), Inches(0.38),
        size=14, bold=True, color=col)
    for j, b in enumerate(bullets):
        by = y + Inches(0.62) + Inches(0.5) * j
        rect(s, x + Inches(0.18), by + Inches(0.12), Inches(0.1), Inches(0.14), col)
        txt(s, b, x + Inches(0.38), by, Inches(3.55), Inches(0.42), size=10, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 10 — Thank You
# ══════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, Inches(0), Inches(0), Inches(0.12), H, TEAL)
rect(s, Inches(0), Inches(6.9), W, Inches(0.6), DARK_CARD)

txt(s, "PropChain", Inches(0.5), Inches(1.4), Inches(12), Inches(1.4),
    size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Blockchain + AI for Transparent Property Registration in India",
    Inches(0.5), Inches(2.85), Inches(12), Inches(0.6),
    size=22, color=TEAL, align=PP_ALIGN.CENTER)

divider(s, Inches(3.6))

summary = [
    ("🔗", "Mock Blockchain",      "10 endpoints, SHA-256 chain, tamper detection"),
    ("🤖", "AI Fraud Detection",   "Textract + Bedrock + 7 rule-based checks"),
    ("📡", "14 REST endpoints",    "Swagger UI /swagger  |  ReDoc /redoc"),
    ("📦", "Full tooling",         "Postman collection  |  run.sh  |  Commands.md"),
]
for i, (icon, title, sub) in enumerate(summary):
    x = Inches(1.0) + Inches(3.0) * i
    txt(s, icon,  x, Inches(3.9),  Inches(2.7), Inches(0.55), size=24, align=PP_ALIGN.CENTER)
    txt(s, title, x, Inches(4.5),  Inches(2.7), Inches(0.42), size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    txt(s, sub,   x, Inches(4.95), Inches(2.7), Inches(0.55), size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

txt(s, "Team OpsAI  |  AI for Bharat Hackathon  |  Powered by AWS",
    Inches(0.5), Inches(7.0), Inches(12), Inches(0.35),
    size=12, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────
OUT = "/Users/gourav/Project/OpsAi/PropChain_Prototype_Progress.pptx"
prs.save(OUT)
print(f"✅  Saved → {OUT}")
print(f"   Slides: {len(prs.slides)}")
